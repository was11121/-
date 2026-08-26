"""本地 AI 图书馆：文档解析、净化、去重、索引、检索和引用。"""

from __future__ import annotations

import csv
import hashlib
import json
import os
import re
import shutil
import uuid
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any


def _now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


@dataclass(slots=True)
class DocumentRecord:
    document_id: str
    title: str
    source: str
    content: str
    content_hash: str
    content_type: str
    status: str = "active"
    created_at: str = ""
    updated_at: str = ""
    tags: list[str] | None = None


class DocumentProcessor:
    SUPPORTED = {".txt", ".md", ".json", ".csv", ".pdf", ".docx", ".pptx", ".xlsx"}

    def extract(self, filename: str, content: bytes | str) -> tuple[str, str]:
        suffix = Path(filename).suffix.lower()
        if suffix not in self.SUPPORTED:
            raise ValueError(f"unsupported document type: {suffix or 'unknown'}")
        raw = content.encode("utf-8") if isinstance(content, str) else content
        if suffix in {".txt", ".md"}:
            return raw.decode("utf-8", errors="replace"), "text/plain"
        if suffix == ".json":
            try:
                return json.dumps(json.loads(raw.decode("utf-8")), ensure_ascii=False, indent=2), "application/json"
            except json.JSONDecodeError:
                return raw.decode("utf-8", errors="replace"), "application/json"
        if suffix == ".csv":
            rows = list(csv.reader(raw.decode("utf-8-sig", errors="replace").splitlines()))
            return "\n".join(" | ".join(row) for row in rows), "text/csv"
        if suffix == ".pdf":
            try:
                from pypdf import PdfReader
                import io

                reader = PdfReader(io.BytesIO(raw))
                pages = [page.extract_text() or "" for page in reader.pages]
                return "\n\n".join(f"[第 {idx + 1} 页]\n{page}" for idx, page in enumerate(pages)), "application/pdf"
            except ImportError:
                raise RuntimeError("PDF support requires pypdf")
        if suffix == ".docx":
            try:
                from docx import Document
                import io

                doc = Document(io.BytesIO(raw))
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip()), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            except ImportError:
                raise RuntimeError("DOCX support requires python-docx")
        if suffix == ".pptx":
            try:
                from pptx import Presentation
                import io

                presentation = Presentation(io.BytesIO(raw))
                texts: list[str] = []
                for number, slide in enumerate(presentation.slides, 1):
                    slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    texts.append(f"[第 {number} 页]\n" + "\n".join(slide_text))
                return "\n\n".join(texts), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            except ImportError:
                raise RuntimeError("PPTX support requires python-pptx")
        if suffix == ".xlsx":
            try:
                from openpyxl import load_workbook
                import io

                book = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
                sheets = []
                for sheet in book.worksheets:
                    rows = [" | ".join("" if cell is None else str(cell) for cell in row) for row in sheet.iter_rows(values_only=True)]
                    sheets.append(f"[工作表: {sheet.title}]\n" + "\n".join(rows))
                return "\n\n".join(sheets), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            except ImportError:
                raise RuntimeError("XLSX support requires openpyxl")
        raise ValueError(f"unsupported document type: {suffix}")


class LocalLibrary:
    def __init__(self, data_dir: str | Path | None = None):
        root = Path(data_dir or os.getenv("MYAGENT_DATA_DIR", Path(__file__).resolve().parents[1] / "data"))
        self.root = root / "library"
        self.root.mkdir(parents=True, exist_ok=True)
        self.index_path = self.root / "index.json"
        self.docs_path = self.root / "documents"
        self.docs_path.mkdir(exist_ok=True)
        self.processor = DocumentProcessor()

    def _load_index(self) -> list[dict[str, Any]]:
        if not self.index_path.exists():
            return []
        try:
            return json.loads(self.index_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            return []

    def _save_index(self, records: list[dict[str, Any]]) -> None:
        self.index_path.write_text(json.dumps(records, ensure_ascii=False, indent=2), encoding="utf-8")

    @staticmethod
    def clean_document(text: str) -> str:
        text = text.replace("\x00", "")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        lines = [line.strip() for line in text.splitlines()]
        seen: set[str] = set()
        output: list[str] = []
        for line in lines:
            if not line:
                if output and output[-1] != "":
                    output.append("")
                continue
            if line in seen and len(line) > 30:
                continue
            seen.add(line)
            output.append(line)
        return "\n".join(output).strip()

    def ingest_document(self, filename: str, content: bytes | str, *, source: str = "upload", tags: list[str] | None = None) -> dict[str, Any]:
        extracted, content_type = self.processor.extract(filename, content)
        cleaned = self.clean_document(extracted)
        if not cleaned:
            raise ValueError("document has no usable text")
        digest = hashlib.sha256(cleaned.encode("utf-8")).hexdigest()
        records = self._load_index()
        existing = next((item for item in records if item.get("content_hash") == digest and item.get("status") == "active"), None)
        if existing:
            return {"status": "duplicate", "document": existing}
        now = _now()
        document = DocumentRecord("doc_" + uuid.uuid4().hex[:12], Path(filename).name, source, cleaned, digest, content_type, created_at=now, updated_at=now, tags=tags or [])
        records.append(asdict(document))
        (self.docs_path / f"{document.document_id}.txt").write_text(cleaned, encoding="utf-8")
        self._save_index(records)
        return {"status": "indexed", "document": asdict(document)}

    def search_library(self, query: str, limit: int = 5) -> list[dict[str, Any]]:
        q_str = (query or "").lower().strip()
        # 提取英文单词/数字和中文字符/二元组分词，提升多语言检索与长句匹配召回率
        eng_tokens = re.findall(r"[a-z0-9]+", q_str)
        cn_chars = re.findall(r"[\u4e00-\u9fff]", q_str)
        cn_bigrams = [cn_chars[i] + cn_chars[i + 1] for i in range(len(cn_chars) - 1)]
        tokens = set(eng_tokens + cn_chars + cn_bigrams)

        scored: list[tuple[float, dict[str, Any]]] = []
        for item in self._load_index():
            if item.get("status") != "active":
                continue
            content = item.get("content", "")
            title = item.get("title", "")
            c_str = content.lower()
            t_str = title.lower()

            if not tokens:
                scored.append((0, {"document_id": item["document_id"], "title": item["title"], "source": item["source"], "snippet": content[:320].replace("\n", " "), "score": 0, "locator": "全文"}))
                continue

            c_eng = re.findall(r"[a-z0-9]+", c_str)
            c_cn = re.findall(r"[\u4e00-\u9fff]", c_str)
            c_bigrams = [c_cn[i] + c_cn[i + 1] for i in range(len(c_cn) - 1)]
            words = set(c_eng + c_cn + c_bigrams)

            overlap = len(tokens & words)
            # 标题命中加权
            title_hit = sum(3 for token in tokens if len(token) >= 2 and token in t_str)
            # 连续短语精准匹配加权
            exact_hit = 5 if any(len(token) >= 2 and token in c_str for token in (eng_tokens + cn_bigrams)) else 0

            total_score = overlap + title_hit + exact_hit
            if total_score > 0:
                snippet = content[:320].replace("\n", " ")
                scored.append((total_score, {"document_id": item["document_id"], "title": item["title"], "source": item["source"], "snippet": snippet, "score": total_score, "locator": "全文"}))

        scored.sort(key=lambda pair: pair[0], reverse=True)
        return [item for _, item in scored[: max(1, min(limit, 20))]]

    def get_citations(self, query: str, limit: int = 5) -> list[dict[str, Any]]:
        return self.search_library(query, limit)

    def archive_library_entry(self, document_id: str) -> bool:
        records = self._load_index()
        changed = False
        for item in records:
            if item.get("document_id") == document_id:
                item["status"] = "archived"
                item["updated_at"] = _now()
                changed = True
        if changed:
            self._save_index(records)
        return changed

    def _owned_by_user(self, item: dict[str, Any], user_id: str) -> bool:
        uid = (user_id or "").strip()
        if not uid:
            return False
        tags = item.get("tags") or []
        if isinstance(tags, str):
            tags = [t.strip() for t in tags.split(",") if t.strip()]
        return uid in tags or f"user:{uid}" in tags or item.get("owner") == uid

    def documents_for_user(self, user_id: str, *, include_untagged: bool = True) -> list[dict[str, Any]]:
        records = self._load_index()
        uid = (user_id or "").strip()
        if not uid:
            return records
        out = []
        for item in records:
            tags = item.get("tags") or []
            untagged = not tags
            if self._owned_by_user(item, uid) or (include_untagged and untagged):
                out.append(item)
        return out

    def delete_documents_for_user(self, user_id: str) -> int:
        """硬删除带有该用户标签的文档。"""
        uid = (user_id or "").strip()
        if not uid:
            return 0
        records = self._load_index()
        keep, removed = [], []
        for item in records:
            if self._owned_by_user(item, uid):
                removed.append(item)
            else:
                keep.append(item)
        if not removed:
            return 0
        self._save_index(keep)
        for item in removed:
            doc_id = item.get("document_id")
            if not doc_id:
                continue
            doc_file = self.docs_path / f"{doc_id}.txt"
            try:
                if doc_file.exists():
                    doc_file.unlink()
            except OSError:
                pass
        return len(removed)

    def delete_library_entry(self, document_id: str) -> bool:
        """硬删除指定文档（Demo 清理使用）。"""
        records = self._load_index()
        keep = [item for item in records if item.get("document_id") != document_id]
        if len(keep) == len(records):
            return False
        self._save_index(keep)
        doc_file = self.docs_path / f"{document_id}.txt"
        try:
            if doc_file.exists():
                doc_file.unlink()
        except OSError:
            pass
        return True

    def correct_library_entry(self, document_id: str, correction: str) -> bool:
        records = self._load_index()
        for item in records:
            if item.get("document_id") == document_id and item.get("status") == "active":
                item["content"] = self.clean_document(correction)
                item["content_hash"] = hashlib.sha256(item["content"].encode("utf-8")).hexdigest()
                item["updated_at"] = _now()
                (self.docs_path / f"{document_id}.txt").write_text(item["content"], encoding="utf-8")
                self._save_index(records)
                return True
        return False
